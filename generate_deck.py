import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set presentation to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color Palette definitions
    BG_COLOR = RGBColor(244, 245, 247)         # #f4f5f7 - Premium light grey/blue
    BRAND_BLUE = RGBColor(78, 98, 125)         # #4e627d - Logo slate/steel blue
    NAVY_DARK = RGBColor(44, 58, 77)           # #2c3a4d - Deep slate/navy for text headers
    TEXT_CHARCOAL = RGBColor(33, 37, 41)       # #212529 - Highly legible soft black
    TEXT_MUTED = RGBColor(108, 122, 137)       # #6c7a89 - Soft gray for subtitles and footers
    WHITE = RGBColor(255, 255, 255)            # For card backgrounds
    CARD_BORDER = RGBColor(220, 224, 230)      # Very light grey border for cards
    
    # Font definitions
    FONT_TITLE = 'Georgia'
    FONT_BODY = 'Arial'
    
    blank_slide_layout = prs.slide_layouts[6]  # Completely blank slide layout
    
    # -------------------------------------------------------------
    # HELPER FUNCTIONS FOR CLEAN & PIXEL-PERFECT LAYOUT
    # -------------------------------------------------------------
    
    def apply_slide_template(slide, slide_num, total_slides=6):
        """Applies solid background and logo to every content slide."""
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Add Logo to Top-Right (Slightly smaller and shifted down for absolute margin safety)
        logo_path = 'logo_transparent.png'
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.4), Inches(0.5), width=Inches(1.15))
            
        # Add subtle footer
        footer_box = slide.shapes.add_textbox(Inches(0.77), Inches(6.9), Inches(11.8), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"Sample Co. Example Solutions  |  The Manifesto  |  Slide {slide_num} of {total_slides}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        
    def add_slide_header(slide, title_text, section_text="MANIFESTO"):
        """Adds a standard structured header with a tiny top-kicker and a main section title."""
        header_box = slide.shapes.add_textbox(Inches(0.77), Inches(0.4), Inches(9.5), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Section kicker
        p0 = tf.paragraphs[0]
        p0.text = section_text.upper()
        p0.font.name = FONT_BODY
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = BRAND_BLUE
        p0.space_after = Pt(2)
        
        # Main Title
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(26)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        
    def draw_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=WHITE, top_bar=True):
        """Draws a professional white card with rounded corners and optionally a brand top bar."""
        # Rounded Rectangle Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        if len(card.adjustments) > 0:
            card.adjustments[0] = 0.04 # Extremely subtle modern roundness
            
        # Top-Accent Bar
        if top_bar:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.12))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND_BLUE
            bar.line.color.rgb = BRAND_BLUE
            
    def draw_left_border_card(slide, left, top, width, height, bg_color=WHITE, border_color=BRAND_BLUE):
        """Draws a white card with a prominent vertical colored accent bar on the left edge."""
        # Rounded Rectangle Card (the container)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        if len(card.adjustments) > 0:
            card.adjustments[0] = 0.04
            
        # Left-Accent Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = border_color
        bar.line.color.rgb = border_color

    # -------------------------------------------------------------
    # SLIDE 1: COVER SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Light solid background
    background1 = slide1.background
    fill1 = background1.fill
    fill1.solid()
    fill1.fore_color.rgb = BG_COLOR
    
    # Large centered card with subtle border
    draw_card(slide1, Inches(1.5), Inches(1.0), Inches(10.33), Inches(5.5), top_bar=True)
    
    # Centered Logo on Cover
    logo_path = 'logo_transparent.png'
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(5.66), Inches(1.8), width=Inches(2.0))
        
    # Title Text Frame (Centered)
    title_box = slide1.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.33), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    
    p_title = tf1.paragraphs[0]
    p_title.text = "THE SAMPLE CO. MANIFESTO"
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY_DARK
    p_title.space_after = Pt(8)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "A Declaration of Unremarkable Intent"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = FONT_TITLE
    p_sub.font.size = Pt(20)
    p_sub.font.italic = True
    p_sub.font.color.rgb = BRAND_BLUE
    p_sub.space_after = Pt(20)
    
    p_foot = tf1.add_paragraph()
    p_foot.text = "Sample Co. Example Solutions  |  EST. 2026"
    p_foot.alignment = PP_ALIGN.CENTER
    p_foot.font.name = FONT_BODY
    p_foot.font.size = Pt(11)
    p_foot.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: THE SANCTITY OF THE PLACEHOLDER (FIXED OVERFLOW & LOGO)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_template(slide2, 2)
    add_slide_header(slide2, "The Sanctity of the Placeholder", "SECTION I")
    
    # Left Column: Elegant Quote Callout Card with Left-Border Accent
    col1_left = Inches(0.77)
    col_width = Inches(5.6)
    card_top = Inches(1.5)
    card_height = Inches(5.1)
    draw_left_border_card(slide2, col1_left, card_top, col_width, card_height)
    
    quote_box = slide2.shapes.add_textbox(col1_left + Inches(0.35), card_top + Inches(0.3), col_width - Inches(0.6), card_height - Inches(0.5))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_top = tf_q.margin_right = tf_q.margin_bottom = 0
    
    pq_header = tf_q.paragraphs[0]
    pq_header.text = "AN UNTAINTED REALITY"
    pq_header.font.name = FONT_BODY
    pq_header.font.size = Pt(10)
    pq_header.font.bold = True
    pq_header.font.color.rgb = BRAND_BLUE
    pq_header.space_after = Pt(10)
    
    pq_quote = tf_q.add_paragraph()
    pq_quote.text = '“\"Your Text Here\" is the most honest sentence ever written in the English language.”'
    # Reduced from 21 to 18 for massive margin safety
    pq_quote.font.name = FONT_TITLE
    pq_quote.font.size = Pt(18)
    pq_quote.font.bold = True
    pq_quote.font.italic = True
    pq_quote.font.color.rgb = NAVY_DARK
    pq_quote.space_after = Pt(14)
    
    pq_body = tf_q.add_paragraph()
    pq_body.text = 'We recognize the placeholder not as an empty vessel, but as a state of pure, unburdened potential. When you look at an "Example Solution," you are looking at perfection before the messy, complicated, and entirely subjective nature of reality ruins it.'
    pq_body.font.name = FONT_BODY
    pq_body.font.size = Pt(10) # Reduced from 10.5 to 10
    pq_body.font.color.rgb = TEXT_CHARCOAL
    pq_body.line_spacing = 1.15
    
    # Right Column: The Three Key Beliefs (Fixed bottom overlap by tuning text density & margins)
    col2_left = Inches(6.97)
    draw_card(slide2, col2_left, card_top, col_width, card_height, top_bar=True)
    
    beliefs_box = slide2.shapes.add_textbox(col2_left + Inches(0.35), card_top + Inches(0.3), col_width - Inches(0.6), card_height - Inches(0.5))
    tf_b = beliefs_box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    
    pb_header = tf_b.paragraphs[0]
    pb_header.text = "PILLARS OF THE PLACEHOLDER"
    pb_header.font.name = FONT_BODY
    pb_header.font.size = Pt(10)
    pb_header.font.bold = True
    pb_header.font.color.rgb = BRAND_BLUE
    pb_header.space_after = Pt(8) # Tightened from 12
    
    # Tuned copy slightly to be more punchy, helping vertical spacing
    beliefs = [
        ("The Blank Slate", "We believe in the beauty of the blank slate that never gets filled, representing pure conceptual harmony."),
        ("The Empty Spreadsheet", "An empty spreadsheet is a pristine work of art, completely untainted by the cruel unpredictability of actual data."),
        ("Occupying Space", "We do not demand attention; we exist purely to occupy a required dimension and provide a calm, reassuring corporate presence.")
    ]
    
    for i, (title, text) in enumerate(beliefs):
        pb_title = tf_b.add_paragraph()
        pb_title.text = f"•  {title}"
        pb_title.font.name = FONT_BODY
        pb_title.font.size = Pt(11) # Reduced from 12 to 11
        pb_title.font.bold = True
        pb_title.font.color.rgb = NAVY_DARK
        pb_title.space_after = Pt(1)
        pb_title.space_before = Pt(4) if i > 0 else Pt(0) # Reduced from 8 to 4
            
        pb_desc = tf_b.add_paragraph()
        pb_desc.text = text
        pb_desc.font.name = FONT_BODY
        pb_desc.font.size = Pt(9.5) # Reduced from 10 to 9.5
        pb_desc.font.color.rgb = TEXT_CHARCOAL
        pb_desc.line_spacing = 1.1

    # -------------------------------------------------------------
    # SLIDE 3: THE THREE PILLARS OF MEDIOCRITY (ADDITIONAL SAFETY)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_template(slide3, 3)
    add_slide_header(slide3, "The Pillars of Example Solutions", "SECTION II")
    
    card_width3 = Inches(3.6)
    card_height3 = Inches(5.1)
    card_top3 = Inches(1.5)
    left_margin3 = Inches(0.77)
    gap3 = Inches(0.5)
    
    pillars = [
        ("THE COMFORT OF THE BELL CURVE", 
         "The 50th Percentile", 
         [
             ("Avoid the Edges", "Operating at the dangerous, unstable edges of the curve demands endless optimization and fear of obsolescence."),
             ("The Safe Middle", "We aim squarely for the absolute average—the mean, the median, and the mode. There is incredible safety in the middle."),
             ("Zero Surprise", "By striving for the average, we guarantee no one will ever be surprised, startled, or offended by what we produce.")
         ]),
        
        ("THE ELEGANCE OF THE BOILERPLATE", 
         "The Template Imperative", 
         [
             ("Bespoke is Waste", "Custom engineering requires thought, time, and effort that is fundamentally a waste of human energy."),
             ("Boilerplate Defense", "Standard terms and conditions have safely protected absolutely nothing of importance for decades. Why change?"),
             ("Template Mandate", "If a template exists, it is a moral imperative to use it exactly as provided without altering fonts or stock photos.")
         ]),
        
        ("PREDICTABILITY AS A VIRTUE", 
         "Structure Without Function", 
         [
             ("No Gimmicks", "Delighting the customer is a trap designed to force companies into a never-ending cycle of escalating gimmicks."),
             ("Pure Structure", "What you see is mathematically precisely what you get. We provide the elegant structure of a product without its function."),
             ("Lowered Hope", "You will never be disappointed because we will never give you a reason to hope for more. Predictability is efficiency.")
         ])
    ]
    
    for idx, (title, sub, bullet_points) in enumerate(pillars):
        c_left = left_margin3 + idx * (card_width3 + gap3)
        draw_card(slide3, c_left, card_top3, card_width3, card_height3, top_bar=True)
        
        box = slide3.shapes.add_textbox(c_left + Inches(0.25), card_top3 + Inches(0.3), card_width3 - Inches(0.5), card_height3 - Inches(0.45))
        tf_c = box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        
        p0 = tf_c.paragraphs[0]
        p0.text = title
        p0.font.name = FONT_BODY
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = BRAND_BLUE
        p0.space_after = Pt(2)
        
        p1 = tf_c.add_paragraph()
        p1.text = sub
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(13) # Reduced from 14 to 13
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(8) # Reduced from 10 to 8
        
        for b_idx, (b_title, b_text) in enumerate(bullet_points):
            p_bullet = tf_c.add_paragraph()
            p_bullet.text = f"•  {b_title}: "
            p_bullet.font.name = FONT_BODY
            p_bullet.font.size = Pt(9.2)  # Reduced from 9.5 to 9.2
            p_bullet.font.bold = True
            p_bullet.font.color.rgb = NAVY_DARK
            p_bullet.space_before = Pt(3) if b_idx > 0 else Pt(0)
            
            run = p_bullet.add_run()
            run.text = b_text
            run.font.name = FONT_BODY
            run.font.size = Pt(9.0)  # Reduced from 9.2 to 9.0
            run.font.bold = False
            run.font.color.rgb = TEXT_CHARCOAL
            p_bullet.line_spacing = 1.1

    # -------------------------------------------------------------
    # SLIDE 4: OUR STANCE AGAINST DISRUPTION (ADDITIONAL SAFETY)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_template(slide4, 4)
    add_slide_header(slide4, "Standing Firm Against \"Disruption\"", "SECTION III")
    
    comp_width = Inches(5.6)
    comp_height = Inches(2.6)
    comp_top = Inches(1.5)
    comp_gap = Inches(0.6)
    
    # Left: Disruption (Hostile)
    draw_card(slide4, left_margin3, comp_top, comp_width, comp_height, top_bar=True)
    box_dis = slide4.shapes.add_textbox(left_margin3 + Inches(0.3), comp_top + Inches(0.3), comp_width - Inches(0.6), comp_height - Inches(0.45))
    tf_dis = box_dis.text_frame
    tf_dis.word_wrap = True
    tf_dis.margin_left = tf_dis.margin_top = tf_dis.margin_right = tf_dis.margin_bottom = 0
    
    pd_title = tf_dis.paragraphs[0]
    pd_title.text = "THE DISRUPTIVE WAY (HOSTILE & AGGRESSIVE)"
    pd_title.font.name = FONT_BODY
    pd_title.font.size = Pt(10)
    pd_title.font.bold = True
    pd_title.font.color.rgb = RGBColor(180, 80, 80)
    pd_title.space_after = Pt(6)
    
    pd_body = tf_dis.add_paragraph()
    pd_body.text = "•  A Toddler's Philosophy: Conglomerates want to 'move fast and break things.' When things break, someone else has to clean them up.\n•  The Cost of Speed: Moving fast leads to sweat, exhaustion, and inevitable errors. It sacrifices peace for frantic movement."
    pd_body.font.name = FONT_BODY
    pd_body.font.size = Pt(9.5) # Reduced from 10 to 9.5
    pd_body.font.color.rgb = TEXT_CHARCOAL
    pd_body.line_spacing = 1.1

    # Right: The Status Quo (Sample Co. Way)
    right_left = left_margin3 + comp_width + comp_gap
    draw_card(slide4, right_left, comp_top, comp_width, comp_height, top_bar=True)
    box_sq = slide4.shapes.add_textbox(right_left + Inches(0.3), comp_top + Inches(0.3), comp_width - Inches(0.6), comp_height - Inches(0.45))
    tf_sq = box_sq.text_frame
    tf_sq.word_wrap = True
    tf_sq.margin_left = tf_sq.margin_top = tf_sq.margin_right = tf_sq.margin_bottom = 0
    
    ps_title = tf_sq.paragraphs[0]
    ps_title.text = "THE SAMPLE CO. WAY (CALM & COMPLIANT)"
    ps_title.font.name = FONT_BODY
    ps_title.font.size = Pt(10)
    ps_title.font.bold = True
    ps_title.font.color.rgb = BRAND_BLUE
    ps_title.space_after = Pt(6)
    
    ps_body = tf_sq.add_paragraph()
    ps_body.text = "•  Moderate Pace: We move at a moderate, legally compliant pace to maintain the status quo.\n•  Anti-Disruption: We view disruption as a hostile act against the comforting, therapeutic rhythm of standard operating procedures. We never reinvent the wheel."
    ps_body.font.name = FONT_BODY
    ps_body.font.size = Pt(9.5) # Reduced from 10 to 9.5
    ps_body.font.color.rgb = TEXT_CHARCOAL
    ps_body.line_spacing = 1.1
    
    # Bottom featured product callout (Sample Wheel Model A)
    bot_top = Inches(4.4)
    bot_width = Inches(11.8)
    bot_height = Inches(2.2)
    
    draw_left_border_card(slide4, left_margin3, bot_top, bot_width, bot_height, border_color=BRAND_BLUE)
    box_wheel = slide4.shapes.add_textbox(left_margin3 + Inches(0.4), bot_top + Inches(0.3), bot_width - Inches(0.8), bot_height - Inches(0.4))
    tf_wh = box_wheel.text_frame
    tf_wh.word_wrap = True
    tf_wh.margin_left = tf_wh.margin_top = tf_wh.margin_right = tf_wh.margin_bottom = 0
    
    pw_title = tf_wh.paragraphs[0]
    pw_title.text = "FEATURED ENGINEERING CASE STUDY"
    pw_title.font.name = FONT_BODY
    pw_title.font.size = Pt(9)
    pw_title.font.bold = True
    pw_title.font.color.rgb = BRAND_BLUE
    pw_title.space_after = Pt(4)
    
    pw_name = tf_wh.add_paragraph()
    pw_name.text = "The 'Sample Wheel' (Model A)"
    pw_name.font.name = FONT_TITLE
    pw_name.font.size = Pt(14) # Reduced from 15 to 14
    pw_name.font.bold = True
    pw_name.font.color.rgb = NAVY_DARK
    pw_name.space_after = Pt(6)
    
    pw_desc = tf_wh.add_paragraph()
    pw_desc.text = "•  No Need to Reinvent: We think the current wheel is perfectly fine (perhaps even a bit too flashy in its roundness).\n•  Sufficient Performance: Model A is available only in gray, functioning at roughly 80% capacity—which is statistically all you really need to get down a slight, predictable incline."
    pw_desc.font.name = FONT_BODY
    pw_desc.font.size = Pt(9.5) # Reduced from 10 to 9.5
    pw_desc.font.color.rgb = TEXT_CHARCOAL
    pw_desc.line_spacing = 1.1

    # -------------------------------------------------------------
    # SLIDE 5: A VISION FOR A GENERIC FUTURE (ADDITIONAL SAFETY)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_template(slide5, 5)
    add_slide_header(slide5, "A Vision for a Generic Future", "SECTION IV")
    
    grid_w = Inches(5.6)
    grid_h = Inches(2.4)
    grid_margin = Inches(0.77)
    grid_gap_x = Inches(0.6)
    grid_gap_y = Inches(0.3)
    
    visions = [
        ("Eradicate Decision Fatigue", 
         "A world where choice is eliminated, completely removing fatigue. Every storefront is named 'Store' and every beverage is labeled 'Drink'."),
        
        ("Absolute Brand Unity", 
         "Every corporate logo is a simple, unassuming letter encased in a mildly rounded square, printed in a calm, non-threatening navy blue."),
        
        ("No More A/B Testing", 
         "There is no A/B testing because there is only Option A. No customized user journeys—just a single, straight, unpaved road from A to B."),
        
        ("The Universal Standard", 
         "John Doe is the CEO of every major corporation. The standard 555 telephone prefix is the only way to communicate.")
    ]
    
    for idx, (title, desc) in enumerate(visions):
        col = idx % 2
        row = idx // 2
        
        v_left = grid_margin + col * (grid_w + grid_gap_x)
        v_top = Inches(1.5) + row * (grid_h + grid_gap_y)
        
        draw_card(slide5, v_left, v_top, grid_w, grid_h, top_bar=True)
        
        v_box = slide5.shapes.add_textbox(v_left + Inches(0.3), v_top + Inches(0.3), grid_w - Inches(0.6), grid_h - Inches(0.45))
        tf_v = v_box.text_frame
        tf_v.word_wrap = True
        tf_v.margin_left = tf_v.margin_top = tf_v.margin_right = tf_v.margin_bottom = 0
        
        pv_title = tf_v.paragraphs[0]
        pv_title.text = f"0{idx+1}.  {title}"
        pv_title.font.name = FONT_TITLE
        pv_title.font.size = Pt(13) # Reduced from 14 to 13
        pv_title.font.bold = True
        pv_title.font.color.rgb = NAVY_DARK
        pv_title.space_after = Pt(4)
        
        pv_desc = tf_v.add_paragraph()
        pv_desc.text = desc
        pv_desc.font.name = FONT_BODY
        pv_desc.font.size = Pt(9.2) # Reduced from 9.5 to 9.2
        pv_desc.font.color.rgb = TEXT_CHARCOAL
        pv_desc.line_spacing = 1.1

    # -------------------------------------------------------------
    # SLIDE 6: CONCLUSION (ADDITIONAL SAFETY)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_slide_template(slide6, 6)
    add_slide_header(slide6, "Join the Unremarkable Movement", "CONCLUSION")
    
    center_w = Inches(10.33)
    center_h = Inches(3.2)
    center_left = Inches(1.5)
    center_top = Inches(1.6)
    
    draw_left_border_card(slide6, center_left, center_top, center_w, center_h, border_color=BRAND_BLUE)
    
    quote_box = slide6.shapes.add_textbox(center_left + Inches(0.5), center_top + Inches(0.4), center_w - Inches(1.0), center_h - Inches(0.6))
    tf_end = quote_box.text_frame
    tf_end.word_wrap = True
    tf_end.margin_left = tf_end.margin_top = tf_end.margin_right = tf_end.margin_bottom = 0
    
    pe_kicker = tf_end.paragraphs[0]
    pe_kicker.text = "THE FINAL IMPERATIVE"
    pe_kicker.font.name = FONT_BODY
    pe_kicker.font.size = Pt(11)
    pe_kicker.font.bold = True
    pe_kicker.font.color.rgb = BRAND_BLUE
    pe_kicker.space_after = Pt(12)
    
    pe_quote = tf_end.add_paragraph()
    pe_quote.text = "“Lower your standards. Embrace the boilerplate. Stop striving for the exceptional and sink comfortably into the warm, inviting embrace of the unremarkable.”"
    pe_quote.font.name = FONT_TITLE
    pe_quote.font.size = Pt(19) # Reduced from 21 to 19 for perfect padding and breathing room
    pe_quote.font.bold = True
    pe_quote.font.italic = True
    pe_quote.font.color.rgb = NAVY_DARK
    pe_quote.space_after = Pt(20)
    
    # Bottom final tagline
    bot_box = slide6.shapes.add_textbox(Inches(1.5), Inches(5.3), Inches(10.33), Inches(1.0))
    tf_bot = bot_box.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    
    pb_text = tf_bot.paragraphs[0]
    pb_text.text = "Sample Co. Example Solutions. We are here. And that is all."
    pb_text.alignment = PP_ALIGN.CENTER
    pb_text.font.name = FONT_TITLE
    pb_text.font.size = Pt(20)
    pb_text.font.bold = True
    pb_text.font.color.rgb = NAVY_DARK
    
    prs.save('Sample_Co_Manifesto.pptx')
    print("Successfully generated Sample_Co_Manifesto.pptx with fixed formatting (iteration 2).")

if __name__ == '__main__':
    create_presentation()
